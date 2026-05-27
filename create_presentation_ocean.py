# -*- coding: utf-8 -*-
"""
生成手势识别项目介绍PPT - Ocean Depths主题优化版
"""

from pptx import Presentation # type: ignore
from pptx.util import Inches, Pt, Cm # type: ignore
from pptx.dml.color import RGBColor # type: ignore
from pptx.enum.text import PP_ALIGN # type: ignore
from pptx.enum.shapes import MSO_SHAPE # type: ignore

# Ocean Depths 主题配色
OCEAN_COLORS = {
    'primary': RGBColor(10, 49, 87),      # 深海蓝
    'secondary': RGBColor(27, 83, 136),   # 海洋蓝
    'accent': RGBColor(64, 156, 214),     # 明亮蓝
    'light': RGBColor(197, 225, 255),     # 浅蓝
    'text': RGBColor(255, 255, 255),      # 白色文字
    'text_dark': RGBColor(30, 50, 70),    # 深色文字
    'background': RGBColor(240, 248, 255) # 背景浅蓝
}

def create_ocean_presentation():
    # 创建演示文稿对象（使用16:9模板）
    prs = Presentation()
    
    # 设置幻灯片大小为16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 获取布局模板
    title_layout = prs.slide_layouts[0]  # 标题页
    content_layout = prs.slide_layouts[1] # 标题+内容
    blank_layout = prs.slide_layouts[5]   # 空白页
    
    # ---------------------- 第1页：封面 ----------------------
    slide = prs.slides.add_slide(title_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['primary']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "实时手势交互识别系统"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['text']
    
    # 添加副标题
    subtitle = slide.placeholders[1]
    subtitle.text = "基于 OpenCV + MediaPipe 的智能手势识别解决方案\n\nGesture Recognition Team | 2026"
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = OCEAN_COLORS['light']
    
    # 添加装饰圆形
    shapes = slide.shapes
    for i in range(3):
        circle = shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(12 + i*1.2), Inches(1 + i*0.4),
                                  Inches(1.5), Inches(1.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = OCEAN_COLORS['accent']
        circle.fill.transparency = 0.3 + i*0.15
        circle.line.fill.background()
    
    # ---------------------- 第2页：目录 ----------------------
    slide = prs.slides.add_slide(content_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['background']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "目录"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['primary']
    
    # 添加目录内容
    content = slide.placeholders[1]
    items = ["项目概述", "功能特性", "技术架构", "快速开始", 
             "手势识别演示", "火球特效展示", "未来展望"]
    
    content.text = ""
    for i, item in enumerate(items):
        para = content.text_frame.add_paragraph()
        para.text = f"{i+1}. {item}"
        para.font.size = Pt(20)
        para.font.color.rgb = OCEAN_COLORS['text_dark']
        para.space_after = Pt(8)
    
    # ---------------------- 第3页：项目概述 ----------------------
    slide = prs.slides.add_slide(content_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['background']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "项目概述"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['primary']
    
    # 添加内容
    content = slide.placeholders[1]
    content.text = """📌 项目背景
手势交互是人机交互的重要方式，无需接触即可控制设备，在智能家居、游戏、医疗等领域有广泛应用。

🎯 项目目标
实时检测手部关键点，识别多种常见手势，提供可视化反馈，并添加趣味特效增强用户体验。"""
    
    for para in content.text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = OCEAN_COLORS['text_dark']
    
    # ---------------------- 第4页：功能特性 ----------------------
    slide = prs.slides.add_slide(content_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['background']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "功能特性"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['primary']
    
    # 添加内容
    content = slide.placeholders[1]
    content.text = """🤚 支持7种手势识别
• 拳头、手掌、点赞
• 剪刀手、OK、食指指上
• 未知手势自动识别

⚡ 实时处理能力
• 摄像头实时视频流
• 毫秒级响应速度
• 流畅的视觉体验

✨ 丰富可视化效果
• 手部关键点标记（21个）
• 骨架连线实时显示
• 酷炫火球特效展示"""
    
    for para in content.text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = OCEAN_COLORS['text_dark']
    
    # ---------------------- 第5页：技术架构 ----------------------
    slide = prs.slides.add_slide(title_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['secondary']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "技术架构"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['text']
    
    # 添加技术栈信息
    subtitle = slide.placeholders[1]
    subtitle.text = """📊 技术栈
OpenCV-Python  —— 视频捕获与图像处理
MediaPipe      —— 手部关键点检测  
PIL/Pillow     —— 中文文本绘制
NumPy          —— 数值计算

🔄 处理流程
摄像头捕获 → 关键点检测 → 手势识别 → 可视化输出"""
    
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = OCEAN_COLORS['light']
    
    # ---------------------- 第6页：快速开始 ----------------------
    slide = prs.slides.add_slide(content_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['background']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "快速开始"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['primary']
    
    # 添加内容
    content = slide.placeholders[1]
    content.text = """📦 安装依赖
$ pip install opencv-python mediapipe Pillow numpy

🚀 运行项目
$ python gesture_recognition.py

⌨️ 操作说明
• 按 'q' 键退出程序
• 按 'r' 键重置帧计数器
• 按 'f' 键切换火球效果"""
    
    for para in content.text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = OCEAN_COLORS['text_dark']
    
    # ---------------------- 第7页：手势识别演示 ----------------------
    slide = prs.slides.add_slide(content_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['background']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "手势识别演示"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['primary']
    
    # 添加内容
    content = slide.placeholders[1]
    content.text = """🎬 支持的手势

👊 拳头        🖐️ 手掌        👍 点赞
✌️ 剪刀手      👌 OK          ☝️ 食指指上

📷 演示特点
• 21个关键点实时追踪
• 绿色骨架连线显示
• 红色圆点标记关键点
• 实时显示手势名称
• 标注左右手信息"""
    
    for para in content.text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = OCEAN_COLORS['text_dark']
    
    # ---------------------- 第8页：火球特效展示 ----------------------
    slide = prs.slides.add_slide(title_layout)
    
    # 设置背景（深色背景突出火球效果）
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 20, 30)
    
    # 添加标题
    title = slide.shapes.title
    title.text = "🔥 火球特效展示"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 150, 50)
    
    # 添加特效特性
    subtitle = slide.placeholders[1]
    subtitle.text = """✨ 特效特性

多层渐变颜色（红→橙→黄→白）
动态粒子系统模拟火焰飘动
跟随手掌实时移动
按 'f' 键可开关控制

💡 使用方式
默认开启火球效果
按 'f' 键切换开关
支持双手同时显示"""
    
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = OCEAN_COLORS['light']
    
    # ---------------------- 第9页：未来展望 ----------------------
    slide = prs.slides.add_slide(content_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['background']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "未来展望"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['primary']
    
    # 添加内容
    content = slide.placeholders[1]
    content.text = """🎯 短期目标
• 增加更多手势类型
• 优化识别准确率
• 支持手势控制应用

🚀 中期目标
• 3D手势识别
• 手势动画效果
• 多手势组合识别

🌟 长期愿景
• 手势控制智能家居
• AR/VR手势交互
• 手语识别与翻译"""
    
    for para in content.text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = OCEAN_COLORS['text_dark']
    
    # ---------------------- 第10页：致谢 ----------------------
    slide = prs.slides.add_slide(title_layout)
    
    # 设置背景
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = OCEAN_COLORS['primary']
    
    # 添加标题
    title = slide.shapes.title
    title.text = "感谢观看！"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = OCEAN_COLORS['text']
    
    # 添加联系方式
    subtitle = slide.placeholders[1]
    subtitle.text = """📧 联系方式: gesture@example.com
🔗 项目地址: github.com/example/gesturerecognition"""
    
    for para in subtitle.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = OCEAN_COLORS['light']
    
    # 保存演示文稿
    prs.save('手势识别项目介绍_Ocean主题.pptx')
    print("PPT文件已生成：手势识别项目介绍_Ocean主题.pptx")

if __name__ == '__main__':
    create_ocean_presentation()