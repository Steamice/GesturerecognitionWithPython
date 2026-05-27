# -*- coding: utf-8 -*-
"""
生成手势识别项目介绍PPT

"""

from pptx import Presentation # type: ignore
from pptx.util import Inches, Pt  # type: ignore
from pptx.dml.color import RGBColor # type: ignore
from pptx.enum.text import PP_ALIGN # type: ignore

def create_presentation():
    # 创建演示文稿对象
    prs = Presentation()
    
    # ---------------------- 第1页：封面 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "实时手势交互识别系统"
    subtitle.text = "基于 OpenCV + MediaPipe 的智能手势识别解决方案\n\n作者: Gesture Recognition Team\n日期: 2026年"
    
    # 设置标题样式
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 39, 97)  # 深蓝色
    
    # 背景渐变
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)  # 浅灰蓝
    
    # ---------------------- 第2页：目录 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "目录"
    
    content = slide.placeholders[1]
    content.text = """
1. 项目概述
2. 功能特性
3. 技术架构
4. 快速开始
5. 手势识别演示
6. 火球特效展示
7. 未来展望
    """
    content.text_frame.paragraphs[0].font.size = Pt(20)
    
    # ---------------------- 第3页：项目概述 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "项目概述"
    
    content = slide.placeholders[1]
    content.text = """
📌 项目背景
- 手势交互是人机交互的重要方式
- 无需接触即可控制设备
- 在智能家居、游戏、医疗等领域有广泛应用

🎯 项目目标
- 实时检测手部关键点
- 识别多种常见手势
- 提供可视化反馈
- 添加趣味特效增强体验
    """
    
    # ---------------------- 第4页：功能特性 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "功能特性"
    
    content = slide.placeholders[1]
    content.text = """
✅ 支持7种手势识别
- 拳头 👊
- 手掌 🖐️
- 点赞 👍
- 剪刀手 ✌️
- OK 👌
- 食指指上 ☝️
- 未知手势

✅ 实时处理
- 摄像头实时视频流
- 毫秒级响应速度
- 流畅的视觉体验

✅ 可视化效果
- 手部关键点标记
- 骨架连线显示
- 手势名称标注
- 酷炫火球特效
    """
    
    # ---------------------- 第5页：技术架构 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "技术架构"
    
    content = slide.placeholders[1]
    content.text = """
📊 技术栈
├── OpenCV-Python (视频捕获与图像处理)
├── MediaPipe (手部关键点检测)
├── PIL/Pillow (中文文本绘制)
└── NumPy (数值计算)

🔄 处理流程
1. 摄像头捕获视频帧
2. MediaPipe 检测关键点
3. 自定义算法识别手势
4. 绘制可视化效果
5. 显示输出画面

🎯 核心算法
- 手指弯曲检测
- 左右手判断
- 手势特征匹配
    """
    
    # ---------------------- 第6页：快速开始 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "快速开始"
    
    content = slide.placeholders[1]
    content.text = """
📦 安装依赖
pip install opencv-python mediapipe Pillow numpy

🚀 运行项目
python gesture_recognition.py

⌨️ 操作说明
- 按 'q' 键退出程序
- 按 'r' 键重置帧计数器
- 按 'f' 键切换火球效果
    """
    
    # ---------------------- 第7页：手势识别演示 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "手势识别演示"
    
    content = slide.placeholders[1]
    content.text = """
🎬 演示效果

手部关键点检测：
- 21个关键点实时追踪
- 绿色骨架连线
- 红色圆点标记

手势识别结果：
- 实时显示手势名称
- 标注左右手信息
- 中文界面友好

📷 摄像头要求
- USB摄像头或内置摄像头
- 建议光线充足环境
- 背景简洁效果更佳
    """
    
    # ---------------------- 第8页：火球特效展示 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "火球特效展示"
    
    content = slide.placeholders[1]
    content.text = """
🔥 火球特效

特效特性：
- 多层渐变颜色（红→橙→黄→白）
- 动态粒子系统
- 跟随手掌移动
- 可开关控制

视觉效果：
- 外层红色光晕
- 中层黄色火焰
- 内层白色核心
- 飘动的火焰粒子

✨ 使用方式
- 默认开启火球效果
- 按 'f' 键切换开关
- 支持双手同时显示
    """
    
    # ---------------------- 第9页：未来展望 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "未来展望"
    
    content = slide.placeholders[1]
    content.text = """
🎯 短期目标
- 增加更多手势类型
- 优化识别准确率
- 支持手势控制应用

🚀 中期目标
- 3D手势识别
- 手势动画效果
- 多手势组合识别

🌟 长期愿景
- 手势控制智能家居
- AR/VR手势交互
- 手语识别与翻译
    """
    
    # ---------------------- 第10页：致谢 ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "感谢观看！"
    
    subtitle = slide.placeholders[1]
    subtitle.text = "实时手势交互识别系统\n\n联系方式：gesture@example.com\n项目地址：github.com/example/gesturerecognition"
    
    # 设置样式
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 39, 97)
    
    # 保存演示文稿
    prs.save('手势识别项目介绍.pptx')
    print("PPT文件已生成：手势识别项目介绍.pptx")

if __name__ == '__main__':
    create_presentation()